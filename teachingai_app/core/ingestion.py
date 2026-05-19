from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterable


MIN_TEXT_THRESHOLD = 80


class ParseError(Exception):
    pass


def _normalize_text(raw_text: str) -> str:
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    return "\n".join(lines)


def _has_enough_text(text: str) -> bool:
    return len(text.strip()) >= MIN_TEXT_THRESHOLD


def _get_paddle_ocr_instance():
    # First try PaddleOCR (original implementation)
    try:
        from paddleocr import PaddleOCR
        try:
            # Use Chinese + English general model for mixed teaching documents.
            return PaddleOCR(use_angle_cls=True, lang="ch")
        except Exception:
            pass
    except Exception:
        pass
    
    # Fallback to easyocr if PaddleOCR is not available
    try:
        import easyocr
        reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)
        return reader
    except Exception:
        return None


def _ocr_image_bytes(image_bytes: bytes, ocr_engine) -> str:
    try:
        from PIL import Image
    except Exception:
        return ""

    if ocr_engine is None:
        return ""

    try:
        image = Image.open(BytesIO(image_bytes)).convert("RGB")
        # Check if it's PaddleOCR or EasyOCR
        if hasattr(ocr_engine, 'ocr'):
            # PaddleOCR
            result = ocr_engine.ocr(image, cls=True)
            chunks: list[str] = []
            for line in result or []:
                for item in line or []:
                    if len(item) >= 2 and isinstance(item[1], (list, tuple)):
                        text = str(item[1][0]).strip()
                        if text:
                            chunks.append(text)
            return "\n".join(chunks)
        else:
            # EasyOCR
            result = ocr_engine.readtext(image)
            chunks: list[str] = []
            for item in result or []:
                if len(item) >= 2:
                    text = str(item[1]).strip()
                    if text:
                        chunks.append(text)
            return "\n".join(chunks)
    except Exception:
        return ""


def _ocr_pdf_pages(content: bytes, ocr_engine) -> str:
    if ocr_engine is None:
        return ""

    try:
        import fitz
    except Exception:
        return ""

    chunks: list[str] = []
    try:
        doc = fitz.open(stream=content, filetype="pdf")
        for page in doc:
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            page_bytes = pix.tobytes("png")
            ocr_text = _ocr_image_bytes(page_bytes, ocr_engine)
            if ocr_text:
                chunks.append(ocr_text)
        doc.close()
    except Exception:
        return ""

    return "\n".join(chunks)


def parse_txt(content: bytes) -> str:
    return _normalize_text(content.decode("utf-8", errors="ignore"))


def parse_docx(content: bytes) -> str:
    try:
        from docx import Document
    except Exception as exc:
        raise ParseError("缺少 python-docx 依赖，无法读取 DOCX 文件。") from exc

    document = Document(BytesIO(content))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    return _normalize_text(text)


def parse_pptx(content: bytes, enable_ocr: bool = True) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ParseError("缺少 python-pptx 依赖，无法读取 PPTX 文件。") from exc

    prs = Presentation(BytesIO(content))
    chunks: list[str] = []
    ocr_engine = _get_paddle_ocr_instance() if enable_ocr else None

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)

            # OCR for image-only slides or pasted screenshots.
            if enable_ocr and hasattr(shape, "shape_type") and getattr(shape, "shape_type", None) == 13:
                image = getattr(shape, "image", None)
                blob = getattr(image, "blob", None)
                if blob:
                    ocr_text = _ocr_image_bytes(blob, ocr_engine)
                    if ocr_text:
                        chunks.append(ocr_text)

    return _normalize_text("\n".join(chunks))


def parse_pdf(content: bytes, enable_ocr: bool = True) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise ParseError("缺少 pypdf 依赖，无法读取 PDF 文件。") from exc

    reader = PdfReader(BytesIO(content))
    chunks: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        chunks.append(text)

    merged = _normalize_text("\n".join(chunks))
    if _has_enough_text(merged) or not enable_ocr:
        return merged

    ocr_engine = _get_paddle_ocr_instance()
    ocr_text = _ocr_pdf_pages(content, ocr_engine)
    if not ocr_text:
        return merged

    return _normalize_text("\n\n".join([merged, ocr_text]))


def parse_file(file_name: str, content: bytes, enable_ocr: bool = True) -> str:
    suffix = Path(file_name).suffix.lower()

    if suffix in {".txt", ".md"}:
        return parse_txt(content)
    if suffix == ".docx":
        return parse_docx(content)
    if suffix == ".pptx":
        return parse_pptx(content, enable_ocr=enable_ocr)
    if suffix == ".pdf":
        return parse_pdf(content, enable_ocr=enable_ocr)

    raise ParseError(f"暂不支持的文件类型: {suffix}")


def merge_text_sources(texts: Iterable[str]) -> str:
    merged = "\n\n".join(t.strip() for t in texts if t and t.strip())
    return _normalize_text(merged)
