from __future__ import annotations

from io import BytesIO
from pathlib import Path

from .models import SimulationReport


class PptxRevisionError(Exception):
    pass


def _pick_slide_layout(prs):
    # Prefer blank layout to avoid placeholder constraints.
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[0]


def _collect_revision_items(report: SimulationReport, limit: int = 12) -> list[str]:
    items: list[str] = []
    seen: set[str] = set()

    def add(text: str) -> None:
        key = text.strip()
        if not key or key in seen:
            return
        seen.add(key)
        items.append(key)

    for item in report.lesson_plan_change_summary:
        add(item)
    for suggestion in report.suggestions:
        add(f"[{suggestion.priority}] {suggestion.suggestion}")
    for record in report.module_deliberations:
        for adj in record.teaching_adjustments:
            add(f"[{record.module_title}] {adj}")

    return items[:limit]


def _append_notes_for_module_adjustments(prs, report: SimulationReport, fallback_slide) -> None:
    if not report.lesson_modules or not report.module_deliberations:
        return

    total_slides = len(prs.slides)
    if total_slides <= 0:
        return

    deliberation_by_module = {d.module_id: d for d in report.module_deliberations}
    for module in sorted(report.lesson_modules, key=lambda m: m.order):
        record = deliberation_by_module.get(module.module_id)
        if record is None or not record.teaching_adjustments:
            continue

        # MVP mapping rule: module order -> slide order (m1->第1页, m2->第2页...)
        slide_index = max(0, module.order - 1)
        target_slide = prs.slides[slide_index] if slide_index < total_slides else fallback_slide

        notes_slide = target_slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        if notes_tf is None:
            continue

        prefix = f"[AI即时调整] 模块{module.order}：{module.title}"
        existing = (notes_tf.text or "").strip()
        if prefix in existing:
            continue

        addition_lines = [prefix]
        addition_lines.extend(
            [f"- {idx}. {item}" for idx, item in enumerate(record.teaching_adjustments, start=1)]
        )
        addition = "\n".join(addition_lines)
        notes_tf.text = f"{existing}\n\n{addition}".strip() if existing else addition


def build_revised_pptx_payload(
    original_file_name: str,
    original_pptx_bytes: bytes,
    report: SimulationReport,
) -> tuple[bytes, str, str]:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as exc:
        raise PptxRevisionError("缺少 python-pptx 依赖，无法生成修订版PPT。") from exc

    try:
        prs = Presentation(BytesIO(original_pptx_bytes))
    except Exception as exc:
        raise PptxRevisionError("原始PPT读取失败，无法生成修订版。") from exc

    slide = prs.slides.add_slide(_pick_slide_layout(prs))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    title_shape = slide.shapes.add_textbox(
        int(slide_w * 0.06),
        int(slide_h * 0.05),
        int(slide_w * 0.88),
        int(slide_h * 0.12),
    )
    title_tf = title_shape.text_frame
    title_tf.text = "AI修订建议（自动生成）"
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.size = Pt(30)

    content_shape = slide.shapes.add_textbox(
        int(slide_w * 0.06),
        int(slide_h * 0.18),
        int(slide_w * 0.88),
        int(slide_h * 0.74),
    )
    tf = content_shape.text_frame
    tf.word_wrap = True

    intro = tf.paragraphs[0]
    intro.text = "说明：本页汇总 AI 对当前课例的可执行修改建议；原始幻灯片内容保持不变。"
    intro.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "重点知识点"
    p.font.bold = True
    p.font.size = Pt(18)

    key_points = report.extracted_key_points[:5] or ["未提取到重点知识点，请结合正文人工确认。"]
    for idx, item in enumerate(key_points, start=1):
        row = tf.add_paragraph()
        row.text = f"{idx}. {item}"
        row.level = 1
        row.font.size = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "建议优先调整项"
    p2.font.bold = True
    p2.font.size = Pt(18)

    revision_items = _collect_revision_items(report)
    if not revision_items:
        revision_items = ["暂无可执行调整项，请查看导出报告中的详细建议。"]
    for idx, item in enumerate(revision_items, start=1):
        row = tf.add_paragraph()
        row.text = f"{idx}. {item}"
        row.level = 1
        row.font.size = Pt(14)

    _append_notes_for_module_adjustments(prs, report, fallback_slide=slide)

    stream = BytesIO()
    prs.save(stream)
    revised_name = f"{Path(original_file_name).stem}_AI修订建议版.pptx"
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return stream.getvalue(), revised_name, mime
